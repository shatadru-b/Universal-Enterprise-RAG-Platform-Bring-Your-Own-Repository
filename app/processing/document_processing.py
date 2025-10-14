# Document normalization, chunking, embedding for multiple file formats

import io
import mimetypes
import re

# Import libraries with graceful fallbacks
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import chardet
except ImportError:
    chardet = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import markdown
except ImportError:
    markdown = None

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False


def detect_encoding(doc_bytes):
    """Detect text encoding using chardet."""
    if chardet:
        result = chardet.detect(doc_bytes)
        return result.get('encoding', 'utf-8')
    return 'utf-8'


def extract_text_from_pdf(doc_bytes):
    """Extract text from PDF using pdfplumber (preferred) or PyPDF2 fallback."""
    text_parts = []
    
    # Try pdfplumber first (better for complex PDFs)
    if pdfplumber:
        try:
            with pdfplumber.open(io.BytesIO(doc_bytes)) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(f"[Page {page_num + 1}]\n{page_text.strip()}")
                            print(f"[PDF] pdfplumber: Page {page_num + 1} extracted {len(page_text)} chars")
                    except Exception as e:
                        print(f"[PDF WARNING] Page {page_num + 1} failed: {e}")
                        
                if text_parts:
                    return "\n\n".join(text_parts)
        except Exception as e:
            print(f"[PDF] pdfplumber failed: {e}, trying PyPDF2...")
    
    # Fallback to PyPDF2
    if PyPDF2:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(doc_bytes))
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        lines = [line.strip() for line in page_text.split('\n') if line.strip()]
                        cleaned_text = '\n'.join(lines)
                        text_parts.append(f"[Page {page_num + 1}]\n{cleaned_text}")
                        print(f"[PDF] PyPDF2: Page {page_num + 1} extracted {len(cleaned_text)} chars")
                except Exception as e:
                    print(f"[PDF WARNING] Page {page_num + 1} failed: {e}")
                    
            if text_parts:
                return "\n\n".join(text_parts)
        except Exception as e:
            print(f"[PDF] PyPDF2 failed: {e}")
    
    return None


def extract_text_from_excel(doc_bytes, filetype):
    """Extract text from Excel files (.xlsx, .xls)."""
    text_parts = []
    
    # Try openpyxl for .xlsx
    if filetype in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"] and openpyxl:
        try:
            wb = openpyxl.load_workbook(io.BytesIO(doc_bytes), data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n[Sheet: {sheet_name}]\n")
                
                # Extract all rows
                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        text_parts.append(" | ".join(row_data))
                        
            if text_parts:
                print(f"[EXCEL] Extracted from {len(wb.sheetnames)} sheets using openpyxl")
                return "\n".join(text_parts)
        except Exception as e:
            print(f"[EXCEL] openpyxl failed: {e}")
    
    # Try pandas for both .xlsx and .xls
    if pd:
        try:
            # Try reading all sheets
            excel_file = pd.ExcelFile(io.BytesIO(doc_bytes))
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_parts.append(f"\n[Sheet: {sheet_name}]\n")
                text_parts.append(df.to_string(index=False))
                
            if text_parts:
                print(f"[EXCEL] Extracted from {len(excel_file.sheet_names)} sheets using pandas")
                return "\n".join(text_parts)
        except Exception as e:
            print(f"[EXCEL] pandas failed: {e}")
    
    # Try xlrd for old .xls files
    if filetype in ["application/vnd.ms-excel", ".xls"] and xlrd:
        try:
            wb = xlrd.open_workbook(file_contents=doc_bytes)
            for sheet in wb.sheets():
                text_parts.append(f"\n[Sheet: {sheet.name}]\n")
                for row_idx in range(sheet.nrows):
                    row_data = [str(cell.value) for cell in sheet.row(row_idx) if cell.value]
                    if row_data:
                        text_parts.append(" | ".join(row_data))
                        
            if text_parts:
                print(f"[EXCEL] Extracted from {wb.nsheets} sheets using xlrd")
                return "\n".join(text_parts)
        except Exception as e:
            print(f"[EXCEL] xlrd failed: {e}")
    
    return None


def extract_text_from_powerpoint(doc_bytes):
    """Extract text from PowerPoint files (.pptx)."""
    if not Presentation:
        return None
        
    try:
        prs = Presentation(io.BytesIO(doc_bytes))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n[Slide {slide_num}]\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
                    
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text])
                        if row_text:
                            text_parts.append(row_text)
        
        if text_parts:
            print(f"[PPTX] Extracted from {len(prs.slides)} slides")
            return "\n".join(text_parts)
    except Exception as e:
        print(f"[PPTX] Extraction failed: {e}")
    
    return None


def extract_text_from_word(doc_bytes):
    """Extract text from Word documents (.docx)."""
    if not DocxDocument:
        return None
        
    try:
        doc = DocxDocument(io.BytesIO(doc_bytes))
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text])
                if row_text:
                    text_parts.append(row_text)
        
        if text_parts:
            print(f"[DOCX] Extracted {len(text_parts)} paragraphs/rows")
            return "\n".join(text_parts)
    except Exception as e:
        print(f"[DOCX] Extraction failed: {e}")
    
    return None


def extract_text_from_csv(doc_bytes):
    """Extract text from CSV files."""
    if not pd:
        # Fallback to basic parsing
        try:
            text = doc_bytes.decode(detect_encoding(doc_bytes))
            return text
        except Exception:
            return None
    
    try:
        df = pd.read_csv(io.BytesIO(doc_bytes))
        text = df.to_string(index=False)
        print(f"[CSV] Extracted {len(df)} rows, {len(df.columns)} columns")
        return text
    except Exception as e:
        print(f"[CSV] pandas failed: {e}, trying basic text extraction")
        try:
            return doc_bytes.decode(detect_encoding(doc_bytes))
        except Exception:
            return None


def extract_text_from_html(doc_bytes):
    """Extract text from HTML files."""
    if not BeautifulSoup:
        # Fallback: return raw text
        try:
            return doc_bytes.decode(detect_encoding(doc_bytes))
        except Exception:
            return None
    
    try:
        soup = BeautifulSoup(doc_bytes, 'lxml')
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator='\n', strip=True)
        print(f"[HTML] Extracted {len(text)} chars")
        return text
    except Exception as e:
        print(f"[HTML] Extraction failed: {e}")
        try:
            return doc_bytes.decode(detect_encoding(doc_bytes))
        except Exception:
            return None


def extract_text_from_image(doc_bytes):
    """Extract text from images using OCR."""
    if not HAS_OCR:
        return "Error: OCR not available. Install Pillow and pytesseract to process images."
    
    try:
        image = Image.open(io.BytesIO(doc_bytes))
        text = pytesseract.image_to_string(image)
        print(f"[OCR] Extracted {len(text)} chars from image")
        return text if text.strip() else "Error: No text found in image."
    except Exception as e:
        print(f"[OCR] Failed: {e}")
        return f"Error: Image OCR failed - {str(e)}"


def normalize_document(doc_bytes, filetype):
    """
    Extracts text from various file formats.
    Supports: PDF, DOCX, XLSX, XLS, PPTX, CSV, TXT, HTML, XML, JSON, MD, RTF, and images.
    """
    print(f"[NORMALIZE] Processing filetype: {filetype}")
    
    # Normalize filetype
    filetype_lower = filetype.lower() if isinstance(filetype, str) else ""
    
    # PDF files
    if any(x in filetype_lower for x in ["pdf", "application/pdf"]):
        result = extract_text_from_pdf(doc_bytes)
        if result:
            return result
        return "Error: Unable to extract text from PDF. The file may be corrupted or image-based."
    
    # Word documents
    elif any(x in filetype_lower for x in ["docx", "wordprocessingml", "msword"]):
        result = extract_text_from_word(doc_bytes)
        if result:
            return result
        return "Error: Unable to extract text from Word document."
    
    # Excel files
    elif any(x in filetype_lower for x in ["xlsx", "xls", "spreadsheetml", "ms-excel"]):
        result = extract_text_from_excel(doc_bytes, filetype_lower)
        if result:
            return result
        return "Error: Unable to extract text from Excel file. Make sure openpyxl or pandas is installed."
    
    # PowerPoint files
    elif any(x in filetype_lower for x in ["pptx", "ppt", "presentationml", "ms-powerpoint"]):
        result = extract_text_from_powerpoint(doc_bytes)
        if result:
            return result
        return "Error: Unable to extract text from PowerPoint file. Make sure python-pptx is installed."
    
    # CSV files
    elif any(x in filetype_lower for x in ["csv", "text/csv"]):
        result = extract_text_from_csv(doc_bytes)
        if result:
            return result
        return "Error: Unable to extract text from CSV file."
    
    # HTML files
    elif any(x in filetype_lower for x in ["html", "htm", "text/html"]):
        result = extract_text_from_html(doc_bytes)
        if result:
            return result
        return "Error: Unable to extract text from HTML file."
    
    # Image files
    elif any(x in filetype_lower for x in ["image/", "png", "jpg", "jpeg", "tiff", "bmp", "gif"]):
        return extract_text_from_image(doc_bytes)
    
    # JSON files
    elif any(x in filetype_lower for x in ["json", "application/json"]):
        try:
            import json
            data = json.loads(doc_bytes.decode(detect_encoding(doc_bytes)))
            text = json.dumps(data, indent=2)
            print(f"[JSON] Extracted {len(text)} chars")
            return text
        except Exception as e:
            print(f"[JSON] Failed: {e}")
            return f"Error: JSON parsing failed - {str(e)}"
    
    # XML files
    elif any(x in filetype_lower for x in ["xml", "application/xml", "text/xml"]):
        result = extract_text_from_html(doc_bytes)  # BeautifulSoup can parse XML too
        if result:
            return result
    
    # Markdown files
    elif any(x in filetype_lower for x in ["markdown", ".md", "text/markdown"]):
        try:
            text = doc_bytes.decode(detect_encoding(doc_bytes))
            if markdown:
                # Convert to HTML first, then extract text
                html = markdown.markdown(text)
                soup = BeautifulSoup(html, 'html.parser') if BeautifulSoup else None
                if soup:
                    return soup.get_text(separator='\n', strip=True)
            return text
        except Exception as e:
            print(f"[MARKDOWN] Failed: {e}")
    
    # Plain text files (default fallback)
    else:
        try:
            encoding = detect_encoding(doc_bytes)
            text = doc_bytes.decode(encoding)
            print(f"[TEXT] Decoded with {encoding} encoding")
            return text
        except Exception as e:
            print(f"[TEXT] Decoding failed: {e}")
            try:
                return doc_bytes.decode(errors="ignore")
            except Exception:
                return "Error: Unable to decode file as text."


def chunk_document(text, chunk_size=512, overlap=64):
    """
    Splits text into overlapping chunks for embedding.
    """
    chunk_size = 1024  # Increased chunk size for better context
    overlap = 128      # Slightly increased overlap
    chunks = []
    start = 0
    import re
    control_re = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]+")
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        # Remove control/binary characters which may appear from bad decoding
        chunk = control_re.sub(" ", chunk)
        # Trim whitespace
        chunk = chunk.strip()
        if not chunk:
            start += chunk_size - overlap
            continue
        chunks.append(chunk)
        print(f"[INGEST DEBUG] Chunk {len(chunks)-1}: {repr(chunk[:100])} ...")
        start += chunk_size - overlap
    print(f"[INGEST DEBUG] Total chunks created: {len(chunks)}")
    return chunks



def embed_chunks(chunks):
    """
    Calls Ollama's embedding API via LangChain for real embeddings.
    """
    # Prefer the official langchain-ollama if available
    try:
        from langchain_ollama import OllamaEmbeddings as OllamaEmbeddingsOfficial
        embedder = OllamaEmbeddingsOfficial(model="all-minilm:22m")
    except Exception:
        try:
            from langchain_community.embeddings import OllamaEmbeddings
            embedder = OllamaEmbeddings(model="all-minilm:22m")
        except Exception:
            raise ImportError("Please install langchain-ollama or langchain-community: pip install langchain-ollama or pip install langchain-community")
    embeddings = embedder.embed_documents(chunks)
    if embeddings and len(embeddings) > 0:
        print(f"[DEBUG] Embedding dimension: {len(embeddings[0])}")
    else:
        print("[DEBUG] No embeddings generated.")
    return embeddings
